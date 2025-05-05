import os
import fitz  # PDF
from pptx import Presentation  # PPTX
from docx import Document as DocxDocument  # DOCX
from PIL import Image
from io import BytesIO
import pytesseract
from typing import List,Dict
from langchain_openai import OpenAIEmbeddings
from config.settings import Config

class ImageTextProcessor:
    def __init__(self):
        self.embeddings = OpenAIEmbeddings(openai_api_key=Config.openai_key)

    def extract_images_from_document(self,file_path: str, output_dir: str) -> list:
        os.makedirs(output_dir, exist_ok=True)
        ext = os.path.splitext(file_path)[1].lower()
        image_paths = []

        if ext == ".pdf":
            doc = fitz.open(file_path)
            for page_num, page in enumerate(doc):
                for img_index, img in enumerate(page.get_images(full=True)):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    image_path = os.path.join(output_dir, f"pdf_page{page_num+1}_img{img_index+1}.{image_ext}")
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    image_paths.append(image_path)

        elif ext == ".pptx":
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                for j, shape in enumerate(slide.shapes):
                    if hasattr(shape, "image"):
                        image = shape.image
                        image_bytes = image.blob
                        image_ext = image.ext
                        image_path = os.path.join(output_dir, f"ppt_slide{i+1}_img{j+1}.{image_ext}")
                        with open(image_path, "wb") as f:
                            f.write(image_bytes)
                        image_paths.append(image_path)

        elif ext == ".docx":
            doc = DocxDocument(file_path)
            for i, shape in enumerate(doc.inline_shapes):
                img_bytes = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
                print("Image extraction from .docx requires `docx2txt` or similar.")

        else:
            raise ValueError("Unsupported document type")

        return image_paths

    def extract_text_from_images(self, image_paths):
        extracted = []
        for path in image_paths:
            image = Image.open(path)
            text = pytesseract.image_to_string(image)
            if text.strip():
                embedding = self.embeddings.embed_query(text)
                extracted.append({
                    "type": "image",
                    "image_path": path,
                    "content": text,
                    "embedding": embedding
                })
        return extracted