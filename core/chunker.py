from typing import List
from langchain_experimental.text_splitter import SemanticChunker
from langchain.text_splitter import RecursiveCharacterTextSplitter,MarkdownHeaderTextSplitter,SentenceTransformersTokenTextSplitter
from .document_loader import DocumentLoader 
from typing import List
from config.settings import Config
from langchain.docstore.document import Document
from pptx import Presentation  

class DocumentChunker:
    def __init__(self, embeddings):
        self.embeddings = embeddings
        self.semantic_chunker = SemanticChunker(
            self.embeddings, 
            min_chunk_size=100
        )
        self.recursive_chunker = RecursiveCharacterTextSplitter(
            chunk_size=Config.CHUNK_SIZE,
            chunk_overlap=Config.CHUNK_OVERLAP
        )
    
    def chunk_pdf(self, file_path: str) -> List[Document]: #Splitting using Semantic Chunker => split by meaning
        docs = DocumentLoader.load_document(file_path)
        splitter = self.semantic_chunker #split doc into chunks, each chunk is of size 600 and each new chunk repeats 100 characters fro the previous one(chunk)=>AI will remeber context better
        chunks = splitter.split_documents(docs)
        return chunks
    
    def chunk_docx(self, file_path: str) -> List[Document]: #Hybrid splitting: semantic(using mMarkDown...) & character splitting if no heading found
        docs = DocumentLoader.load_document(file_path)
        
        # Combine text for Markdown heading detection
        raw_text = "\n".join([doc.page_content for doc in docs])
        heading_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=[("#", "title"), ("##", "section")])
        heading_chunks = heading_splitter.split_text(raw_text)

        if len(heading_chunks) <= 1: #means we have probably any clear headings => semantic splitting
            print("No heading detected — falling back to Semantic Chunker")
            splitter = self.semantic_chunker
            chunks = splitter.split_documents(docs)
        else:
            chunks = heading_chunks

        return chunks

    def chunk_txt(self, file_path: str) -> List[Document]:#Splitting using Semantic Chunker => split by meaning
        with open(file_path, "r") as f:
            raw_text = f.read()
            docs = [Document(page_content=raw_text)]
            splitter = self.semantic_chunker
            chunks = splitter.split_documents(docs)
        return chunks

    def chunk_pptx(self, file_path: str) -> List[Document]:
        prs = Presentation(file_path)
        slide_docs = []
        
        for i, slide in enumerate(prs.slides, start=1):
            slide_content = []
            
            # Extract slide title if available
            if slide.shapes.title and slide.shapes.title.text:
                slide_content.append(f"# {slide.shapes.title.text.strip()}")
            
            # Process all shapes on the slide
            for shape in slide.shapes:
                # Text from text boxes, tables, etc.
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                
                # Handle tables specially
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_data.append(cell.text.strip())
                        if row_data:
                            table_data.append(" | ".join(row_data))
                    
                    if table_data:
                        slide_content.append("TABLE:\n" + "\n".join(table_data))
                
                # Notes about charts and images
                if shape.has_chart:
                    slide_content.append("[CHART: Chart data not extractable as text]")
                
                # Extract alt text from images when available
                if hasattr(shape, "image") and hasattr(shape, "alt_text") and shape.alt_text:
                    slide_content.append(f"[IMAGE: {shape.alt_text}]")
            
            # Get notes if available
            notes_text = ""
            if (slide.has_notes_slide and slide.notes_slide.notes_text_frame and slide.notes_slide.notes_text_frame.text.strip()):
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                slide_content.append(f"NOTES: {notes_text}")
            
            # Combine all slide content
            full_text = "\n\n".join(slide_content)
            
            if full_text:  # Only create document if slide has content
                doc = Document(
                    page_content=full_text,
                    metadata={
                        "source": file_path,
                        "slide_number": i,
                        "total_slides": len(prs.slides),
                        "has_notes": bool(notes_text),
                        "file_type": "pptx"
                    }
                )
                slide_docs.append(doc)
        
        return slide_docs
        
    def chunk_document(self, file_path: str, doc_type: str) -> List[Document]:
        if doc_type == ".pdf":
            return self.chunk_pdf(file_path)
        elif doc_type == '.pptx':
            return self.chunk_pptx(file_path)
        elif doc_type == '.docx':
            return self.chunk_docx(file_path)
        elif doc_type == '.txt':
            return self.chunk_txt(file_path)
        else:
            raise ValueError(f"Unsupported document type: {doc_type}")
        
